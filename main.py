from pptx import Presentation
from pptx.shapes.placeholder import LayoutPlaceholder, PlaceholderPicture
import pandas as pd
import copy
from math import ceil
import os

#region Input parameters
DATA_FILE_PATH = "./IUI25_1badgeXslide/people.csv" # The CSV file containing the data to use for the badges
TEMPLATE_PPTX_PATH = "./IUI25/badge-layout.pptx" # The template file to use for the badges
OUTPUT_PPTX_PATH = "./IUI25/OUTPUT.pptx" # The output file will be saved here. If the file already exists, it will be overwritten. If the file does not exist, it will be created.
HARD_CODED_DATA = True # If True, the data is hardcoded in the script. If False, the data is read from a CSV file

# Placeholders in the PowerPoint template
PLACEHOLDER_NAME = "Nome " # The placeholder for the name
PLACEHOLDER_SURNAME = "COGNOME" # The placeholder for the surname
PLACEHOLDER_UNIVERSITY = "Workshops + Main Conference" # The placeholder for the university
N_TEMPLATE_PLACEHOLDERS = 3 # The number of placeholders in the template (for each badge)
SLD_LAYOUT = 6 # The layout of the slide to duplicate (0-indexed)
N_BADGE_PER_SLIDE = 1 # The number of badges per slide

# Data keys in the CSV file
K_NAME = "name" # The key in the data for the name
K_SURNAME = "surname" # The key in the data for the surname
K_UNIVERSITY = "type" # The key in the data for the university or any 3rd field you want to add to the badge
#endregion

def get_data():
    # return a pandas dataframe with 80 random rows. Each row represent a person at a conference. Each person as a name, a surname, a University.
    if HARD_CODED_DATA:
        data = {
            "name": ["John", "Jane", "Alice", "Bob", "Charlie", "David", "Eve", "Frank", "Grace", "Hannah", "Ivan", "Jack", "Katie", "Liam", "Mia", "Nathan", "Olivia", "Peter", "Quinn", "Rachel", "Steve", "Tina", "Umberto", "Violet", "Walter", "Xavier", "Yvonne", "Zach", "Albert", "Bella", "Carmen", "Dylan", "Ella", "Fiona", "George", "Helen", "Igor", "Jenny", "Karl", "Lara", "Micheal", "Nora", "Oscar", "Pamela", "Quentin", "Rita", "Sam", "Tara", "Ugo", "Valeria", "William", "Xena", "Yuri", "Zoe", "Alessandro", "Beatrice", "Carlo", "Davide", "Elena", "Fabio", "Giulia", "Hugo", "Irene", "Jorge", "Klara", "Luca", "Maria", "Nico", "Ottavia", "Paolo", "Quirino", "Riccardo", "Sara", "Tommaso", "Umberto", "Valentina", "Walter", "Xena", "Yuri", "Zoe"],
            "surname": ["Doe", "Smith", "Brown", "Johnson", "Williams", "Jones", "Garcia", "Martinez", "Hernandez", "Lopez", "Gonzalez", "Wilson", "Anderson", "Thomas", "Taylor", "Moore", "Jackson", "White", "Harris", "Martin", "Thompson", "Garcia", "Martinez", "Hernandez", "Lopez", "Gonzalez", "Wilson", "Anderson", "Thomas", "Taylor", "Moore", "Jackson", "White", "Harris", "Martin", "Thompson", "Garcia", "Martinez", "Hernandez", "Lopez", "Gonzalez", "Wilson", "Anderson", "Thomas", "Taylor", "Moore", "Jackson", "White", "Harris", "Martin", "Thompson", "Garcia", "Martinez", "Hernandez", "Lopez", "Gonzalez", "Wilson", "Anderson", "Thomas", "Taylor", "Moore", "Jackson", "White", "Harris", "Martin", "Thompson", "Garcia", "Martinez", "Hernandez", "Lopez", "Gonzalez", "Wilson", "Anderson", "Thomas", "Taylor", "Moore", "Jackson", "White", "Harris", "Martin"],
            "type": ["University of Banana", "University of Apple", "University of Pear", "University of Orange", "University of Strawberry", "University of Cherry", "University of Kiwi", "University of Pineapple", "University of Lemon", "University of Blueberry", "University of Raspberry", "University of Blackberry", "University of Watermelon", "University of Melon", "University of Grape", "University of Peach", "University of Plum", "University of Apricot", "University of Mango", "University of Papaya", "University of Guava", "University of Pomegranate", "University of Coconut", "University of Kiwi", "University of Pineapple", "University of Lemon", "University of Blueberry", "University of Raspberry", "University of Blackberry", "University of Watermelon", "University of Melon", "University of Grape", "University of Peach", "University of Plum", "University of Apricot", "University of Mango", "University of Papaya", "University of Guava", "University of Pomegranate", "University of Coconut", "University of Kiwi", "University of Pineapple", "University of Lemon", "University of Blueberry", "University of Raspberry", "University of Blackberry", "University of Watermelon", "University of Melon", "University of Grape", "University of Peach", "University of Plum", "University of Apricot", "University of Mango", "University of Papaya", "University of Guava", "University of Pomegranate", "University of Coconut", "University of Kiwi", "University of Pineapple", "University of Lemon", "University of Blueberry", "University of Raspberry", "University of Blackberry", "University of Watermelon", "University of Melon", "University of Grape", "University of Peach", "University of Plum", "University of Apricot", "University of Mango", "University of Papaya", "University of Guava", "University of Pomegranate", "University of Coconut", "University of Kiwi", "University of Pineapple", "University of Lemon", "University of Blueberry", "University of Atheneum", "Unversity of Lyceum"]
        }
        return pd.DataFrame(data)
    else:
        df = pd.read_csv(DATA_FILE_PATH)
        return df  # Return all rows

# def duplicate_slide(pres, slide_index):
#     template = prs.slide_layouts[slide_index]
#     new_slide = pres.slides.add_slide(template)
#
#     # Copy the elements from the template slide to the new slide
#     for shape in template.shapes:
#         if shape.has_text_frame:
#             new_shape = new_slide.shapes.add_shape(
#                 shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height
#             )
#             new_shape.text = shape.text_frame.text
#         elif shape.has_chart:
#             # For charts, you might need to copy data and properties explicitly
#             # Here we're keeping it simple and skipping charts
#             pass
#         else:
#             # For other types of shapes, you may need to handle them similarly
#             pass
#
#     return new_slide

def copy_slide_from_external_prs(prs):

    # copy from external presentation all objects into the existing presentation
    external_pres = Presentation(TEMPLATE_PPTX_PATH)

    # specify the slide you want to copy the contents from
    ext_slide = external_pres.slides[0]

    # Define the layout you want to use from your generated pptx
    slide_layout = prs.slide_layouts[SLD_LAYOUT]

    # create now slide, to copy contents to 
    curr_slide = prs.slides.add_slide(slide_layout)

    # now copy contents from external slide, but do not copy slide properties
    # e.g. slide layouts, etc., because these would produce errors, as diplicate
    # entries might be generated

    for shp in ext_slide.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return prs


def clean_string(s: str, keep_only_first_word: bool = False) -> str:
    """
    Cleans a given string by applying the following transformations:
    - Converts to title case (first letter uppercase, rest lowercase).
    - If multiple words exist, makes all of them title case.
    - Optionally keeps only the first word if keep_only_first_word is True.
    - Strips leading and trailing whitespace.
    """
    s = s.strip()  # Remove leading and trailing spaces

    if keep_only_first_word:
        s = s.split()[0] if s else ""  # Keep only the first word

    return s.title()

if __name__ == '__main__':
    data = get_data()
    n = len(data)
    
    print(f"Dataframe with {n} rows")

    N_SLIDE_TO_DUPLICATE = ceil((n-N_BADGE_PER_SLIDE) / N_BADGE_PER_SLIDE)

    # The powerpoint should have only 1 slide with the template.
    prs = Presentation(TEMPLATE_PPTX_PATH)

    # Duplicate the slide as many times as needed
    # e.g. 2 badge for slide, 10 people, 5 slides, but 1 already exists, so you will need to duplicate the template 4 times
    for _ in range(N_SLIDE_TO_DUPLICATE):
        prs = copy_slide_from_external_prs(prs)

    idx_person=0
    c=0 # In this template, there are 3 placeholders, so I need to increment the index of the person only when I have changed all the placeholders. c is a counter for this purpose

    # For each slide
    for slide in prs.slides:
        # Check all the elements (a Shape can be nested in a group of shapes, so I need to check all the shapes in the slide)
        for shape in slide.shapes:
            # If there is not text, it's a group of shapes (old template)
            if not shape.has_text_frame:
                print("No text frame")

                if type(shape) == PlaceholderPicture:
                    continue

                #### OLD TEMPLATE ###
                # Per ogni shape del gruppo (che è un badge)
                if shape.shapes:
                    print("\tBut shapes")
                    # Per ogni elemento del badge
                    for s in shape.shapes:
                        # Se c'è un elemento di testo
                        if s.has_text_frame:
                            for paragraph in s.text_frame.paragraphs:
                                # Controllo se ho finito le persone (capita con un numero dispari di persone se i badge sono 2 per slide)
                                if idx_person >= n:
                                    break
                                # Prendi i dati della persona
                                person = data.iloc[idx_person]

                                # Un paragraph può avere più run
                                for run in paragraph.runs:
                                    print("\t\t Before|",run.text,'|')
                                    # Se è un placeholder, cambia il testo con i dati della persona
                                    if  PLACEHOLDER_NAME == run.text:
                                        run.text = person["name"]
                                    elif PLACEHOLDER_SURNAME == run.text:
                                        run.text = person["surname"]
                                    elif PLACEHOLDER_UNIVERSITY in run.text:
                                        run.text = person["university"]
                                    print("\t\t After|",run.text,'|')
                # Finito un badge, passa alla persona successiva
                idx_person += 1
                ### END OLD TEMPLATE ###
            else:
                print("Has text frame")
                s = shape # Alias for the shape

                # Check all the paragraphs in the text_frame
                for paragraph in s.text_frame.paragraphs:
                    # Controllo se ho finito le persone (capita con un numero dispari di persone se i badge sono 2 per slide)
                    if idx_person >= n:
                        break

                    # Prendi i dati della persona corrente
                    person = data.iloc[idx_person]

                    # Un paragraph può avere più run. Cosa sia un run non l'ho ancora capito. Forse una frase del paragrafo? Comunque un sottoinsieme del paragrafo
                    for run in paragraph.runs:
                        print("\t\t Before|",run.text,'|')

                        # Se è un placeholder, cambia il testo con i dati della persona
                        if  PLACEHOLDER_NAME == run.text:
                            # run.text = person[K_NAME].strip()
                            run.text = clean_string(person[K_NAME])
                            c+=1
                        elif PLACEHOLDER_SURNAME == run.text:
                            # run.text = person[K_SURNAME].strip()
                            run.text = clean_string(person[K_SURNAME])
                            c+=1
                        elif PLACEHOLDER_UNIVERSITY in run.text:
                            # run.text = person[K_UNIVERSITY].strip()
                            run.text = clean_string(person[K_UNIVERSITY])
                            c+=1

                        print("\t\t After|",run.text,'|')
                # It checked all the paragraphs in this shape
                if c == N_TEMPLATE_PLACEHOLDERS:
                    idx_person += 1
                    c=0

    # Ensure the output file exists before saving
    if not os.path.exists(OUTPUT_PPTX_PATH):
        open(OUTPUT_PPTX_PATH, 'w').close()
    prs.save(OUTPUT_PPTX_PATH)
